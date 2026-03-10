import subprocess
import tempfile
import os
import json
import re
import logging

logger = logging.getLogger("vietlearn.ocr")

CLAUDE_PATH = "/home/claude-user/.local/bin/claude"
CLAUDE_ENV = {**os.environ, "PATH": "/home/claude-user/.local/bin:" + os.environ.get("PATH", "")}

PROMPT = """Analyse ce contenu. C'est une lecon ou liste de vocabulaire vietnamien.

Extrais TOUS les mots et phrases de vocabulaire que tu trouves.
Pour chaque mot/phrase, donne la traduction francaise.

Reponds UNIQUEMENT avec du JSON valide, sans markdown, sans ```json, juste le JSON brut.
Format exact:
{
  "entries": [
    {"vietnamese": "mot en vietnamien", "french": "traduction francaise", "category": "categorie"},
    ...
  ],
  "description": "description courte du contenu"
}

Pour la categorie, deduis-la du contexte (ex: "Salutations", "Nourriture", "Nombres", "Verbes", "Phrases courantes", etc.)

Si tu vois des phrases completes, inclus-les aussi.
Si tu vois des tons ou de la phonetique, inclus-les dans le champ vietnamese.
Extrais TOUT ce qui est pertinent pour apprendre le vietnamien."""

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
TEXT_DOC_EXTENSIONS = {".doc", ".docx", ".odt", ".rtf", ".ppt", ".pptx", ".xls", ".xlsx", ".pdf", ".txt"}


def extract_vocab_with_ai(file_bytes: bytes, filename: str) -> dict:
    """Analyze any file and extract Vietnamese vocabulary using Claude."""
    suffix = os.path.splitext(filename)[1].lower()

    if suffix in IMAGE_EXTENSIONS:
        return _process_image(file_bytes, suffix)
    elif suffix in TEXT_DOC_EXTENSIONS:
        return _process_text_document(file_bytes, suffix)
    else:
        return _error("Format non supporte: " + suffix)


# --- Images: send directly to Claude (native vision) ---

def _process_image(image_bytes: bytes, suffix: str) -> dict:
    """Send image directly to Claude — he has native vision."""
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(image_bytes)
        tmp_path = tmp.name

    try:
        return _call_claude_with_image(tmp_path)
    finally:
        _safe_delete(tmp_path)


def _call_claude_with_image(image_path: str) -> dict:
    """Call Claude CLI with an image file."""
    cmd = [CLAUDE_PATH, "--print", PROMPT, image_path]
    return _run_claude(cmd)


# --- Documents: extract text, then send to Claude ---

def _process_text_document(file_bytes: bytes, suffix: str) -> dict:
    """Extract text from document, split into chunks, process each with Claude."""
    text = _extract_text(file_bytes, suffix)

    if not text or not text.strip():
        return _error(f"Aucun texte extrait du fichier {suffix}")

    logger.info("Extracted %d chars from %s", len(text), suffix)

    # Split into chunks by page markers or by size
    chunks = _split_into_chunks(text, max_chars=3000)
    logger.info("Split into %d chunks", len(chunks))

    all_entries = []
    descriptions = []

    for i, chunk in enumerate(chunks):
        logger.info("Processing chunk %d/%d (%d chars)", i + 1, len(chunks), len(chunk))
        result = _call_claude_with_text(chunk)
        entries = result.get("entries", [])
        all_entries.extend(entries)
        desc = result.get("description", "")
        if desc:
            descriptions.append(desc)
        logger.info("Chunk %d: %d entries", i + 1, len(entries))

    # Deduplicate by vietnamese text
    seen = set()
    unique_entries = []
    for e in all_entries:
        key = e.get("vietnamese", "").strip().lower()
        if key and key not in seen:
            seen.add(key)
            unique_entries.append(e)

    logger.info("Total: %d entries (%d unique)", len(all_entries), len(unique_entries))

    return {
        "entries": unique_entries,
        "description": " | ".join(descriptions) if descriptions else "",
        "method": "ai",
        "pages": len(chunks),
        "debug_raw": f"{len(chunks)} chunks, {len(unique_entries)} mots extraits",
    }


def _split_into_chunks(text: str, max_chars: int = 3000) -> list[str]:
    """Split text into chunks, respecting page markers or paragraph breaks."""
    # First try splitting by page markers
    pages = re.split(r'---\s*(?:Page|Slide)\s+\d+\s*---', text)
    pages = [p.strip() for p in pages if p.strip()]

    if len(pages) > 1:
        # Merge small pages together up to max_chars
        chunks = []
        current = ""
        for page in pages:
            if len(current) + len(page) > max_chars and current:
                chunks.append(current)
                current = page
            else:
                current = (current + "\n\n" + page).strip()
        if current:
            chunks.append(current)
        return chunks

    # No page markers: split by paragraphs
    if len(text) <= max_chars:
        return [text]

    chunks = []
    paragraphs = text.split("\n")
    current = ""
    for para in paragraphs:
        if len(current) + len(para) > max_chars and current:
            chunks.append(current)
            current = para
        else:
            current = (current + "\n" + para).strip()
    if current:
        chunks.append(current)

    return chunks if chunks else [text]


def _extract_text(file_bytes: bytes, suffix: str) -> str:
    """Extract raw text from various document formats."""
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        if suffix == ".pdf":
            return _text_from_pdf(tmp_path)
        elif suffix in (".docx", ".doc"):
            return _text_from_docx(tmp_path)
        elif suffix in (".pptx", ".ppt"):
            return _text_from_pptx(tmp_path)
        elif suffix in (".xlsx", ".xls"):
            return _text_from_xlsx(tmp_path)
        elif suffix == ".txt":
            return file_bytes.decode("utf-8", errors="replace")
        elif suffix in (".odt", ".rtf"):
            return _text_from_libreoffice(tmp_path)
        else:
            return ""
    finally:
        _safe_delete(tmp_path)


def _text_from_pdf(path: str) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"--- Page {i+1} ---\n{text}")
    return "\n\n".join(pages)


def _text_from_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _text_from_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        lines.append(para.text)
    return "\n".join(lines)


def _text_from_xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"--- {sheet} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines)


def _text_from_libreoffice(path: str) -> str:
    """Convert ODT/RTF to text via LibreOffice."""
    tmpdir = tempfile.mkdtemp(prefix="vietlearn_lo_")
    try:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmpdir, path],
            capture_output=True, timeout=30,
            env={**os.environ, "HOME": tmpdir}
        )
        txt_files = [f for f in os.listdir(tmpdir) if f.endswith(".txt")]
        if txt_files:
            with open(os.path.join(tmpdir, txt_files[0]), "r", errors="replace") as f:
                return f.read()
        return ""
    finally:
        for f in os.listdir(tmpdir):
            _safe_delete(os.path.join(tmpdir, f))
        try:
            os.rmdir(tmpdir)
        except OSError:
            pass


def _call_claude_with_text(text: str) -> dict:
    """Call Claude CLI with text content via stdin."""
    full_prompt = PROMPT + "\n\nVoici le contenu du document:\n\n" + text
    cmd = [CLAUDE_PATH, "--print", full_prompt]
    return _run_claude(cmd)


# --- Common ---

def _run_claude(cmd: list[str]) -> dict:
    """Execute Claude CLI and parse response."""
    try:
        logger.info("Running: %s", " ".join(cmd[:3]) + " ...")
        result = subprocess.run(
            cmd, capture_output=True, text=True,
            timeout=180, env=CLAUDE_ENV
        )

        raw = result.stdout.strip()
        stderr = result.stderr.strip()

        logger.info("Claude returned %d chars (exit %d)", len(raw), result.returncode)
        if stderr:
            logger.warning("Claude stderr: %s", stderr[:300])

        parsed = _parse_json(raw)

        return {
            "entries": parsed.get("entries", []),
            "description": parsed.get("description", ""),
            "method": "ai",
            "pages": 0,
            "debug_raw": raw[:2000],
        }

    except subprocess.TimeoutExpired:
        logger.error("Claude timeout")
        return _error("Timeout — fichier trop volumineux")
    except Exception as e:
        logger.error("Claude error: %s", e)
        return _error(str(e))


def _parse_json(text: str) -> dict:
    """Parse JSON from Claude's response."""
    # Strip markdown fences
    text = re.sub(r'^```json\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'^```\s*$', '', text, flags=re.MULTILINE)
    text = text.strip()

    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # Find JSON object in text
    match = re.search(r'\{[\s\S]*\}', text)
    if match:
        try:
            return json.loads(match.group())
        except json.JSONDecodeError:
            pass

    logger.error("JSON parse failed. Raw: %s", text[:500])
    return {"entries": [], "description": "Erreur de parsing IA"}


def _error(msg: str) -> dict:
    return {
        "entries": [],
        "description": msg,
        "method": "error",
        "pages": 0,
        "debug_raw": msg,
    }


def _safe_delete(path: str):
    try:
        os.unlink(path)
    except OSError:
        pass
